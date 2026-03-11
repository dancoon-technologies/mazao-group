"""
Generate Mazao Group System Tutorial PowerPoint.
Run: pip install python-pptx && python scripts/generate_tutorial_ppt.py
Output: docs/Mazao_Group_System_Tutorial.pptx
"""

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("Install python-pptx first: pip install python-pptx")
    raise

# Output path
OUT_DIR = Path(__file__).resolve().parent.parent / "docs"
OUT_DIR.mkdir(exist_ok=True)
OUT_FILE = OUT_DIR / "Mazao_Group_System_Tutorial.pptx"

# Theme colors (green/organization style)
TITLE_COLOR = RGBColor(0x1B, 0x5E, 0x20)   # Dark green
BULLET_COLOR = RGBColor(0x2E, 0x7D, 0x32)  # Green
BODY_COLOR = RGBColor(0x21, 0x21, 0x21)    # Dark gray


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    left, top, width, height = Inches(0.5), Inches(1.2), Inches(9), Inches(1.2)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    if subtitle:
        left, top, width, height = Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
        tx2 = slide.shapes.add_textbox(left, top, width, height)
        tx2.text_frame.text = subtitle
        tx2.text_frame.paragraphs[0].font.size = Pt(18)
        tx2.text_frame.paragraphs[0].font.color.rgb = BODY_COLOR
    return slide


def add_section_slide(prs, title, bullets, sub_bullets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left, top, width, height = Inches(0.5), Inches(0.4), Inches(9), Inches(0.8)
    tx = slide.shapes.add_textbox(left, top, width, height)
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    left, top, width, height = Inches(0.5), Inches(1.2), Inches(9), Inches(5.5)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.add_paragraph()
        p.text = "• " + line
        p.font.size = Pt(16)
        p.font.color.rgb = BODY_COLOR
        p.space_after = Pt(6)
        if sub_bullets and i < len(sub_bullets) and sub_bullets[i]:
            for sub in sub_bullets[i]:
                p2 = tf.add_paragraph()
                p2.text = "  ◦ " + sub
                p2.font.size = Pt(14)
                p2.font.color.rgb = BODY_COLOR
                p2.space_after = Pt(2)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1. Title
    add_title_slide(
        prs,
        "Mazao Group System Tutorial",
        "Extension Officer Visit Monitoring — How to use all features and roles",
    )

    # 2. System overview
    add_section_slide(
        prs,
        "System Overview",
        [
            "Purpose: Monitor extension officers and prevent ghost visits using GPS + photo verification.",
            "Platforms:",
            "Key concepts: Farmers (people), Farms (plots), Visits (recorded with proof), Schedules (planned visits).",
        ],
        [
            [],
            ["Web app — Supervisors and Admins (dashboards, visits list, schedule approval, staff management).", "Mobile app — Extension Officers (record visits, add farmers, propose schedules)."],
            [],
        ],
    )

    # 3. User roles
    add_section_slide(
        prs,
        "User Roles",
        [
            "Admin — Full access: all data, staff management, invite/deactivate users, assign department/region.",
            "Supervisor — Department/region scope: dashboard, visits list, approve/reject schedule proposals. No Staff page.",
            "Extension Officer — Own data only: assigned farmers, own visits and schedules. Can record visits, add farmers, propose schedules.",
        ],
    )

    # 4. Where to use what
    add_section_slide(
        prs,
        "Where to Use the System",
        [
            "Web app: Used by Admin and Supervisor (and Officers for farmers/schedules only).",
            "Mobile app: Used mainly by Extension Officers; Admins/Supervisors can also use it (e.g. to create schedules and assign officers).",
            "Farmers do not log in — they are managed by staff in the system.",
        ],
    )

    # 5. Logging in (Web)
    add_section_slide(
        prs,
        "Logging In (Web)",
        [
            "Open the web app and go to the login page.",
            "Enter your email and password, then sign in.",
            "If you must change password: you will be redirected to Change Password; set a new password and continue.",
            "After login: Admins and Supervisors go to Dashboard; Officers go to Farmers.",
        ],
    )

    # 6. Logging in (Mobile)
    add_section_slide(
        prs,
        "Logging In (Mobile)",
        [
            "Open the mobile app and enter email and password.",
            "If required, complete Change Password first.",
            "Use Unlock (e.g. biometric) when reopening the app to access the main tabs.",
            "Tabs: Home, Visits, Record (center button), Farmers, Profile.",
        ],
    )

    # 7. Dashboard (Admin & Supervisor)
    add_section_slide(
        prs,
        "Dashboard (Admin & Supervisor only)",
        [
            "View organisation-wide (Admin) or department-scoped (Supervisor) stats and charts.",
            "Stats: visits today/month, active officers, and similar metrics.",
            "Charts: visits by day, visits by activity type, top officers, schedules summary.",
            "Use the day-range selector (e.g. 7, 14, 30 days) for time-based charts.",
        ],
    )

    # 8. Farmers & Farms (Web)
    add_section_slide(
        prs,
        "Farmers & Farms (Web)",
        [
            "Farmers: List and create farmers; assign an Extension Officer to each farmer.",
            "Farms: List and create farms (plot/location); link each farm to a farmer (county, sub-county, village, location, plot size, crop).",
            "Officers see only their assigned farmers; Admins/Supervisors see farmers in scope (department/region or all).",
        ],
    )

    # 9. Farmers & Farms (Mobile)
    add_section_slide(
        prs,
        "Farmers & Farms (Mobile)",
        [
            "Farmers tab: View your assigned farmers and their details.",
            "Add farmer: Use the flow to add a new farmer (name, phone) and at least one farm (location, county, sub-county, village, plot size, crop).",
            "You can add more farms to an existing farmer from farmer detail.",
        ],
    )

    # 10. Visits (Web — Admin & Supervisor)
    add_section_slide(
        prs,
        "Visits (Web — Admin & Supervisor only)",
        [
            "Visits page: List all visits in scope; filter by officer and date.",
            "View visit details: farmer, farm (if any), officer, GPS, photo, activity type, verification status, notes.",
            "Verify: Visits are verified or rejected by the system using GPS distance and required photo; you review and monitor status.",
        ],
    )

    # 11. Recording a Visit (Mobile)
    add_section_slide(
        prs,
        "Recording a Visit (Mobile)",
        [
            "Tap the center Record button (or go to Record Visit from the app).",
            "Select the farmer (or add a new farmer first). Optionally select a specific farm (plot).",
            "Capture photo and allow GPS; add activity type, notes, and any report fields as required.",
            "Submit. The system checks distance from your location to the farmer/farm; visits beyond the threshold are rejected (anti–ghost visit).",
        ],
    )

    # 12. Schedules (Web)
    add_section_slide(
        prs,
        "Schedules (Web)",
        [
            "All roles can open the Schedules page.",
            "Admin/Supervisor: Create schedules (date, optional farmer, notes) and assign an officer; view and approve or reject proposals.",
            "Officer: View own schedules and proposals only.",
            "Approval: Supervisors (and Admins) accept or reject schedule proposals from officers; officers can see status.",
        ],
    )

    # 13. Schedules (Mobile — Propose & Edit)
    add_section_slide(
        prs,
        "Schedules (Mobile — Propose & Edit)",
        [
            "Propose schedule: Create a new schedule (date, optional farmer, notes). Officers propose for themselves; Admins/Supervisors can assign an officer.",
            "Edit schedule: Open a schedule and edit date, farmer, or notes; save changes.",
            "After a Supervisor/Admin approves or rejects, the status updates; check notifications for updates.",
        ],
    )

    # 14. Staff (Admin only)
    add_section_slide(
        prs,
        "Staff (Admin only)",
        [
            "Staff page: List all staff; invite new users (email, role: Supervisor or Extension Officer).",
            "Edit staff: Change department, region, or role; deactivate a user (removes access, keeps data).",
            "Resend credentials: Send password-set or login instructions again to a user.",
        ],
    )

    # 15. Notifications
    add_section_slide(
        prs,
        "Notifications",
        [
            "Web: Bell icon shows unread count; open to list notifications, mark as read, or archive.",
            "Mobile: List on the Notifications screen; tap to open. Push notifications can open this screen.",
            "Use for schedule approvals, important updates, and system messages.",
        ],
    )

    # 16. Locations & options
    add_section_slide(
        prs,
        "Locations & Options",
        [
            "Locations: Use Kenya regions/counties (and sub-county, village) when adding farms or filtering.",
            "Options: Departments (Mazao na afya, Agritech, AgriPriize), staff roles, and activity types are loaded from the backend for dropdowns and filters.",
        ],
    )

    # 17. Quick reference by role
    add_section_slide(
        prs,
        "Quick Reference by Role",
        [
            "Admin: Dashboard, Farmers, Farms, Visits, Schedules, Staff; invite/edit/deactivate users; approve schedules.",
            "Supervisor: Dashboard, Farmers, Farms, Visits, Schedules; approve/reject schedule proposals; no Staff.",
            "Extension Officer (web): Farmers, Farms, Schedules (own); no Dashboard, no Visits list, no Staff.",
            "Extension Officer (mobile): Home, Visits, Record visit, Farmers, Propose/Edit schedule, Notifications, Profile.",
        ],
    )

    # 18. Summary
    add_section_slide(
        prs,
        "Summary",
        [
            "Mazao Group system helps prevent ghost visits with GPS + photo verification.",
            "Three roles: Admin (full), Supervisor (department-scoped), Extension Officer (own data).",
            "Web: dashboards, visit lists, schedule approval, staff management (admin). Mobile: record visits, add farmers, propose/edit schedules.",
            "For API and deployment details, see backend/docs/ and the project README.",
        ],
    )

    prs.save(OUT_FILE)
    print(f"Saved: {OUT_FILE}")


if __name__ == "__main__":
    main()
